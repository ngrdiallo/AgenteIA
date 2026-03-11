"""
Document loaders: estrae testo + metadata da PDF, DOCX, PPTX, XLSX, immagini.
PORTED FROM: AgenteIA-Production/src/pdf_parser.py (esteso a multi-formato)
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}


@dataclass
class DocumentPage:
    """Una singola pagina/sezione di un documento con metadati."""
    text: str
    page_number: int
    source_file: str
    file_type: str
    total_pages: int = 1
    tables: List[List[List[str]]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


class DocumentLoader:
    """
    Loader multi-formato con fallback OCR per PDF scansionati.
    Formati supportati: PDF, DOCX, PPTX, XLSX, immagini (OCR).
    """

    def __init__(self, ocr_enabled: bool = True):
        self.ocr_enabled = ocr_enabled

    def load(self, file_path: str | Path) -> List[DocumentPage]:
        """Auto-detect formato e carica documento."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File non trovato: {path}")

        ext = path.suffix.lower()
        loaders = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".pptx": self._load_pptx,
            ".xlsx": self._load_xlsx,
        }
        image_exts = {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}

        if ext in loaders:
            return loaders[ext](path)
        elif ext in image_exts:
            return self._load_image(path)
        else:
            raise ValueError(f"Formato non supportato: {ext}")

    # ------------------------------------------------------------------
    # PDF: PyMuPDF nativo + OCR fallback per pagine scansionate
    # ------------------------------------------------------------------
    def _load_pdf(self, path: Path) -> List[DocumentPage]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF non installato, fallback a PyPDF2")
            return self._load_pdf_pypdf2(path)

        pages: List[DocumentPage] = []
        doc = fitz.open(str(path))
        total = len(doc)

        for page_num in range(total):
            page = doc[page_num]
            text = page.get_text("text")

            # Se la pagina ha poco testo, prova OCR
            if len(text.strip()) < 50 and self.ocr_enabled:
                ocr_text = self._ocr_page(page)
                if ocr_text:
                    text = ocr_text

            # Estrai tabelle (heuristiche PyMuPDF)
            tables = []
            try:
                tabs = page.find_tables()
                for tab in tabs:
                    tables.append(tab.extract())
            except Exception:
                pass

            pages.append(DocumentPage(
                text=text.strip(),
                page_number=page_num + 1,
                source_file=path.name,
                file_type="pdf",
                total_pages=total,
                tables=tables,
            ))

        doc.close()
        logger.info(f"PDF caricato: {path.name} ({total} pagine)")
        return pages

    def _load_pdf_pypdf2(self, path: Path) -> List[DocumentPage]:
        """Fallback se PyMuPDF non disponibile."""
        import PyPDF2
        pages: List[DocumentPage] = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            total = len(reader.pages)
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                pages.append(DocumentPage(
                    text=text.strip(),
                    page_number=i + 1,
                    source_file=path.name,
                    file_type="pdf",
                    total_pages=total,
                ))
        logger.info(f"PDF (PyPDF2) caricato: {path.name} ({total} pagine)")
        return pages

    def _ocr_page(self, page) -> Optional[str]:
        """OCR su una pagina PDF scansionata via pytesseract."""
        try:
            from PIL import Image
            import pytesseract

            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img, lang="ita+eng")
            return text.strip() if text.strip() else None
        except ImportError:
            logger.debug("pytesseract non installato, OCR non disponibile")
            return None
        except Exception as e:
            logger.debug(f"OCR fallito: {e}")
            return None

    # ------------------------------------------------------------------
    # DOCX: python-docx
    # ------------------------------------------------------------------
    def _load_docx(self, path: Path) -> List[DocumentPage]:
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Estrai tabelle
        tables = []
        for table in doc.tables:
            t = []
            for row in table.rows:
                t.append([cell.text for cell in row.cells])
            tables.append(t)

        full_text = "\n".join(paragraphs)
        # Dividi in pagine virtuali da ~3000 caratteri
        pages = self._split_virtual_pages(full_text, path.name, "docx", tables=tables)
        logger.info(f"DOCX caricato: {path.name} ({len(pages)} pagine)")
        return pages

    # ------------------------------------------------------------------
    # PPTX: python-pptx
    # ------------------------------------------------------------------
    def _load_pptx(self, path: Path) -> List[DocumentPage]:
        from pptx import Presentation

        prs = Presentation(str(path))
        pages: List[DocumentPage] = []
        total = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                # Tabelle nelle slide
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        texts.append(row_text)

            pages.append(DocumentPage(
                text="\n".join(texts),
                page_number=slide_num,
                source_file=path.name,
                file_type="pptx",
                total_pages=total,
            ))

        logger.info(f"PPTX caricato: {path.name} ({total} slide)")
        return pages

    # ------------------------------------------------------------------
    # XLSX: openpyxl
    # ------------------------------------------------------------------
    def _load_xlsx(self, path: Path) -> List[DocumentPage]:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        pages: List[DocumentPage] = []

        for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip(" |"):
                    rows.append(row_text)

            table_data = [
                [str(cell) if cell is not None else "" for cell in row]
                for row in ws.iter_rows(values_only=True)
            ]

            pages.append(DocumentPage(
                text=f"--- Foglio: {sheet_name} ---\n" + "\n".join(rows),
                page_number=sheet_num,
                source_file=path.name,
                file_type="xlsx",
                total_pages=len(wb.sheetnames),
                tables=[table_data] if table_data else [],
            ))

        wb.close()
        logger.info(f"XLSX caricato: {path.name} ({len(pages)} fogli)")
        return pages

    # ------------------------------------------------------------------
    # Immagini: OCR + metadata
    # ------------------------------------------------------------------
    def _load_image(self, path: Path) -> List[DocumentPage]:
        text = ""
        if self.ocr_enabled:
            try:
                from PIL import Image
                import pytesseract

                img = Image.open(str(path))
                text = pytesseract.image_to_string(img, lang="ita+eng")
            except ImportError:
                logger.debug("pytesseract non installato")
            except Exception as e:
                logger.debug(f"OCR immagine fallito: {e}")

        if not text.strip():
            text = f"[Immagine: {path.name} — testo non estraibile via OCR]"

        page = DocumentPage(
            text=text.strip(),
            page_number=1,
            source_file=path.name,
            file_type=path.suffix.lstrip("."),
            total_pages=1,
            metadata={"image_path": str(path)},
        )
        logger.info(f"Immagine caricata: {path.name}")
        return [page]

    # ------------------------------------------------------------------
    # Helper: pagine virtuali per documenti senza paginazione nativa
    # ------------------------------------------------------------------
    def _split_virtual_pages(
        self,
        text: str,
        source_file: str,
        file_type: str,
        page_size: int = 3000,
        tables: Optional[List] = None,
    ) -> List[DocumentPage]:
        """Divide testo lungo in pagine virtuali."""
        if len(text) <= page_size:
            return [DocumentPage(
                text=text,
                page_number=1,
                source_file=source_file,
                file_type=file_type,
                total_pages=1,
                tables=tables or [],
            )]

        pages: List[DocumentPage] = []
        start = 0
        page_num = 1
        while start < len(text):
            end = start + page_size
            # Cerca fine frase più vicina
            if end < len(text):
                for sep in [". ", "\n\n", "\n", " "]:
                    pos = text.rfind(sep, start, end + 200)
                    if pos > start:
                        end = pos + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk:
                pages.append(DocumentPage(
                    text=chunk,
                    page_number=page_num,
                    source_file=source_file,
                    file_type=file_type,
                    total_pages=0,  # calcolato dopo
                    tables=tables if page_num == 1 else [],
                ))
                page_num += 1
            start = end

        # Aggiorna total_pages
        for p in pages:
            p.total_pages = len(pages)

        return pages
